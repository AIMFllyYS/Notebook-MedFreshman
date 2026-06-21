"""
PPTX/PPT → Markdown 降级解析脚本（含图片提取）

用法：
    python scripts/fallback-pptx.py <input.pptx> <output.md> [--subject physics]
    python scripts/fallback-pptx.py --batch --subject physics --dir "path/to/ppts"

依赖安装：
    pip install python-pptx
"""

import sys
import os
import hashlib
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError:
    print("ERROR: python-pptx not installed. Run: pip install python-pptx")
    sys.exit(1)

ROOT = Path(__file__).resolve().parent.parent


def pptx_to_markdown(pptx_path: str, output_path: str, subject: str = "physics") -> int:
    """Convert PPTX to markdown with image extraction. Returns image count."""
    prs = Presentation(pptx_path)
    md_lines: list[str] = []
    basename = Path(pptx_path).stem
    md_lines.append(f"# {basename}\n")

    # Setup image output directory
    img_dir = ROOT / "public" / "images" / subject / basename
    img_dir.mkdir(parents=True, exist_ok=True)
    img_count = 0

    for i, slide in enumerate(prs.slides, 1):
        md_lines.append(f"\n## Slide {i}\n")

        for shape in slide.shapes:
            # Text extraction
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if not text:
                        continue
                    if shape == slide.shapes.title:
                        md_lines.append(f"### {text}\n")
                    else:
                        md_lines.append(f"{text}\n")

            # Image extraction
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                img_bytes = image.blob
                ext = image.content_type.split("/")[-1]
                if ext == "jpeg":
                    ext = "jpg"

                # Hash-based filename for dedup
                img_hash = hashlib.sha256(img_bytes).hexdigest()
                img_filename = f"{img_hash}.{ext}"
                img_filepath = img_dir / img_filename

                if not img_filepath.exists():
                    img_filepath.write_bytes(img_bytes)

                img_path = f"/images/{subject}/{basename}/{img_filename}"
                md_lines.append(f"\n![](/images/{subject}/{basename}/{img_filename})\n")
                img_count += 1

            # Table extraction
            if shape.has_table:
                table = shape.table
                rows = list(table.rows)
                if rows:
                    header = "| " + " | ".join(
                        cell.text.strip() for cell in rows[0].cells
                    ) + " |"
                    separator = "| " + " | ".join("---" for _ in rows[0].cells) + " |"
                    md_lines.append(header)
                    md_lines.append(separator)
                    for row in rows[1:]:
                        row_text = "| " + " | ".join(
                            cell.text.strip() for cell in row.cells
                        ) + " |"
                        md_lines.append(row_text)
                    md_lines.append("")

        # Notes
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                md_lines.append(f"\n> 备注：{notes}\n")

    os.makedirs(os.path.dirname(os.path.abspath(output_path)), exist_ok=True)
    with open(output_path, "w", encoding="utf-8") as f:
        f.write("\n".join(md_lines))

    print(f"  OK {Path(pptx_path).name} -> {output_path}")
    if img_count > 0:
        print(f"    {img_count} images -> public/images/{subject}/{basename}/")

    return img_count


def batch_process(directory: str, subject: str = "physics"):
    """Process all .pptx and .ppt files in a directory."""
    dir_path = Path(directory)
    output_dir = ROOT / "content" / "_raw" / subject / "courseware"
    output_dir.mkdir(parents=True, exist_ok=True)

    files = sorted(list(dir_path.glob("*.pptx")) + list(dir_path.glob("*.ppt")))
    if not files:
        print(f"No .pptx/.ppt files found in {directory}")
        return

    print(f"=== PPT Image Recovery ({len(files)} files) ===\n")
    total_images = 0
    success = 0
    failed = []

    for f in files:
        output_path = output_dir / f"{f.stem}.md"
        try:
            imgs = pptx_to_markdown(str(f), str(output_path), subject)
            total_images += imgs
            success += 1
        except Exception as e:
            print(f"  x {f.name}: {e}")
            failed.append(f.name)

    print(f"\n=== Done ===")
    print(f"  Success: {success}/{len(files)}")
    print(f"  Images recovered: {total_images}")
    if failed:
        print(f"  Failed: {', '.join(failed)}")


if __name__ == "__main__":
    if "--batch" in sys.argv:
        subject = "physics"
        directory = ""
        args = sys.argv[1:]
        i = 0
        while i < len(args):
            if args[i] == "--subject" and i + 1 < len(args):
                subject = args[i + 1]
                i += 2
            elif args[i] == "--dir" and i + 1 < len(args):
                directory = args[i + 1]
                i += 2
            else:
                i += 1
        if not directory:
            print("Usage: python scripts/fallback-pptx.py --batch --subject physics --dir <path>")
            sys.exit(1)
        batch_process(directory, subject)
    elif len(sys.argv) >= 3:
        subject = "physics"
        if "--subject" in sys.argv:
            idx = sys.argv.index("--subject")
            subject = sys.argv[idx + 1]
        pptx_to_markdown(sys.argv[1], sys.argv[2], subject)
    else:
        print("Usage:")
        print("  python scripts/fallback-pptx.py <input.pptx> <output.md> [--subject physics]")
        print("  python scripts/fallback-pptx.py --batch --subject physics --dir <path>")
        sys.exit(1)
